#!/usr/bin/env python3
"""mcp_office_server.py — Fase 20 — MCP server `anja_office`.

Genera DOCX/XLSX/PPTX da spec JSON. Conversion via pandoc/libreoffice/marp.

Tool esposti:
  - office.generate_docx(spec, out_path) — Word document
  - office.generate_xlsx(spec, out_path) — Excel spreadsheet
  - office.generate_pptx(spec, out_path) — PowerPoint slides
  - office.from_markdown(md, target_format, out_path) — pandoc md→docx/pptx
  - office.to_pdf(source_path, out_path?) — libreoffice headless → PDF
  - office.from_markdown_slides(md, theme, out_path) — marp-cli → pptx/pdf
  - office.deps_status() — diagnostica deps installate

Path: out_path è relativo allo scope (hub o workspace), whitelist files/.

Sandbox: stesso pattern di mcp_code_server (env scrub, cwd whitelist).
"""

import json
import os
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Optional


PROTO_VERSION = "2024-11-05"
SERVER_NAME = "anja_office"
SERVER_VERSION = "0.1.0"

SCOPE = os.environ.get("ANJA_SCOPE", "hub")
ROOT = Path(os.environ.get("ANJA_ROOT", os.getcwd())).resolve()


def _hub_root() -> Optional[Path]:
    if SCOPE == "hub":
        return ROOT
    env_hub = os.environ.get("ANJA_HUB")
    if env_hub:
        return Path(env_hub).expanduser().resolve()
    return None


def _resolve_scope_root(scope_arg: Optional[str]) -> tuple:
    """Risolve root del scope per output paths. scope_arg può essere 'hub' o 'workspace:<name>'."""
    s = (scope_arg or SCOPE).strip()
    hub = _hub_root()
    if s == "hub":
        if not hub:
            return None, "hub not resolvable"
        return hub, None
    if s.startswith("workspace:"):
        if not hub:
            return None, "hub not resolvable for workspace scope"
        name = s.split(":", 1)[1].strip()
        ws = hub / "workspaces" / name
        if ws.is_symlink():
            ws = ws.resolve()
        if not ws.is_dir():
            return None, f"workspace '{name}' not found"
        if (ws / ".anjawiki").is_dir():
            return ws / ".anjawiki", None
        return ws, None
    return None, f"invalid scope: {s}"


def _validate_out_path(scope_root: Path, rel_path: str) -> tuple:
    """Path traversal guard + whitelist files/."""
    if not rel_path:
        return None, "out_path required"
    rel = rel_path.lstrip("/").lstrip("\\")
    if ".." in rel:
        return None, "path traversal not allowed"
    parts = [p for p in rel.split("/") if p]
    if not parts or parts[0] != "files":
        return None, "out_path must start with files/"
    target = (scope_root / rel).resolve()
    try:
        target.relative_to(scope_root.resolve())
    except ValueError:
        return None, "path outside scope"
    return target, None


# ============================================================
# DOCX
# ============================================================

def tool_generate_docx(args: dict) -> dict:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx"}

    spec = args.get("spec")
    out_path = args.get("out_path", "")
    scope_arg = args.get("scope")
    if not isinstance(spec, dict):
        return {"error": "spec must be dict"}

    scope_root, err = _resolve_scope_root(scope_arg)
    if err:
        return {"error": err}
    target, err = _validate_out_path(scope_root, out_path)
    if err:
        return {"error": err}

    doc = Document()
    if spec.get("title"):
        doc.add_heading(spec["title"], level=0)

    for section in spec.get("sections", []):
        if "heading" in section:
            doc.add_heading(section["heading"], level=int(section.get("level", 1)))
        if section.get("content"):
            doc.add_paragraph(section["content"])
        if section.get("bullets"):
            for b in section["bullets"]:
                doc.add_paragraph(str(b), style="List Bullet")
        if section.get("table"):
            rows = section["table"]
            if rows:
                tbl = doc.add_table(rows=len(rows), cols=len(rows[0]))
                tbl.style = "Light Grid"
                for i, row in enumerate(rows):
                    for j, cell in enumerate(row):
                        tbl.cell(i, j).text = str(cell)
        if section.get("image"):
            img_path = (scope_root / section["image"]).resolve()
            try:
                img_path.relative_to(scope_root.resolve())
                if img_path.is_file():
                    doc.add_picture(str(img_path))
                    if section.get("caption"):
                        doc.add_paragraph(section["caption"], style="Caption")
            except (ValueError, Exception):
                pass

    target.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(target))
    return {
        "ok": True,
        "format": "docx",
        "path": out_path,
        "absolute": str(target),
        "size": target.stat().st_size,
    }


# ============================================================
# XLSX
# ============================================================

def tool_generate_xlsx(args: dict) -> dict:
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    except ImportError:
        return {"error": "openpyxl not installed. Run: pip install openpyxl"}

    spec = args.get("spec")
    out_path = args.get("out_path", "")
    scope_arg = args.get("scope")
    if not isinstance(spec, dict):
        return {"error": "spec must be dict"}

    scope_root, err = _resolve_scope_root(scope_arg)
    if err:
        return {"error": err}
    target, err = _validate_out_path(scope_root, out_path)
    if err:
        return {"error": err}

    wb = Workbook()
    default_sheet = wb.active
    sheets_spec = spec.get("sheets", [])
    if not sheets_spec:
        return {"error": "spec.sheets required (list of {name, rows, charts?})"}

    first = True
    for sheet_spec in sheets_spec:
        name = (sheet_spec.get("name") or "Sheet")[:31]
        if first:
            ws = default_sheet
            ws.title = name
            first = False
        else:
            ws = wb.create_sheet(title=name)

        rows = sheet_spec.get("rows", [])
        for r_idx, row in enumerate(rows, start=1):
            for c_idx, val in enumerate(row, start=1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        # Charts
        for ch_spec in sheet_spec.get("charts", []):
            ch_type = ch_spec.get("type", "bar")
            data_range = ch_spec.get("data_range")  # es. "A1:B10"
            if not data_range:
                continue
            try:
                start_cell, end_cell = data_range.split(":")
                from openpyxl.utils.cell import coordinate_from_string, column_index_from_string
                col_s, row_s = coordinate_from_string(start_cell)
                col_e, row_e = coordinate_from_string(end_cell)
                ref = Reference(ws, min_col=column_index_from_string(col_s), min_row=row_s,
                                max_col=column_index_from_string(col_e), max_row=row_e)
                if ch_type == "line":
                    chart = LineChart()
                elif ch_type == "pie":
                    chart = PieChart()
                else:
                    chart = BarChart()
                chart.add_data(ref, titles_from_data=bool(ch_spec.get("titles_from_data", True)))
                chart.title = ch_spec.get("title", "")
                anchor = ch_spec.get("anchor", "E2")
                ws.add_chart(chart, anchor)
            except Exception as e:
                print(f"[anja_office] chart skipped: {e}", file=sys.stderr)

    target.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(target))
    return {
        "ok": True,
        "format": "xlsx",
        "path": out_path,
        "absolute": str(target),
        "size": target.stat().st_size,
    }


# ============================================================
# PPTX
# ============================================================

def tool_generate_pptx(args: dict) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    spec = args.get("spec")
    out_path = args.get("out_path", "")
    scope_arg = args.get("scope")
    if not isinstance(spec, dict):
        return {"error": "spec must be dict"}

    scope_root, err = _resolve_scope_root(scope_arg)
    if err:
        return {"error": err}
    target, err = _validate_out_path(scope_root, out_path)
    if err:
        return {"error": err}

    prs = Presentation()
    for slide_spec in spec.get("slides", []):
        layout_name = slide_spec.get("layout", "title_content")
        layout_idx = {"title": 0, "title_content": 1, "content": 5, "blank": 6}.get(layout_name, 1)
        try:
            layout = prs.slide_layouts[layout_idx]
        except IndexError:
            layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Title
        if slide_spec.get("title") and slide.shapes.title:
            slide.shapes.title.text = slide_spec["title"]

        # Bullets in content placeholder
        if slide_spec.get("bullets"):
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # content
                    tf = ph.text_frame
                    tf.clear()
                    for i, b in enumerate(slide_spec["bullets"]):
                        if i == 0:
                            tf.text = str(b)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(b)
                    break

        # Subtitle
        if slide_spec.get("subtitle") and len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = slide_spec["subtitle"]
            except Exception:
                pass

        # Image
        if slide_spec.get("image"):
            img_path = (scope_root / slide_spec["image"]).resolve()
            try:
                img_path.relative_to(scope_root.resolve())
                if img_path.is_file():
                    left = Inches(slide_spec.get("img_left", 1))
                    top = Inches(slide_spec.get("img_top", 1.5))
                    width = Inches(slide_spec.get("img_width", 8))
                    slide.shapes.add_picture(str(img_path), left, top, width=width)
            except (ValueError, Exception):
                pass

    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target))
    return {
        "ok": True,
        "format": "pptx",
        "path": out_path,
        "absolute": str(target),
        "size": target.stat().st_size,
    }


# ============================================================
# Conversion tools (CLI subprocess)
# ============================================================

def _run_cli(cmd: list, timeout: int = 120) -> dict:
    """Esegue CLI con timeout. Ritorna {ok, stdout, stderr, exit_code}."""
    try:
        proc = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        return {
            "ok": proc.returncode == 0,
            "exit_code": proc.returncode,
            "stdout": proc.stdout[-2000:] if proc.stdout else "",
            "stderr": proc.stderr[-2000:] if proc.stderr else "",
        }
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": f"timeout after {timeout}s"}
    except FileNotFoundError:
        return {"ok": False, "error": f"binary not found: {cmd[0]}"}
    except Exception as e:
        return {"ok": False, "error": f"{type(e).__name__}: {e}"}


def tool_from_markdown(args: dict) -> dict:
    """pandoc: markdown → docx/pptx."""
    if not shutil.which("pandoc"):
        return {"error": "pandoc not installed. brew install pandoc"}
    md = args.get("md", "")
    target_format = args.get("target_format", "docx").lower()
    out_path = args.get("out_path", "")
    scope_arg = args.get("scope")
    if not md.strip():
        return {"error": "md required"}
    if target_format not in ("docx", "pptx", "html", "pdf", "odt"):
        return {"error": f"unsupported target_format: {target_format}"}

    scope_root, err = _resolve_scope_root(scope_arg)
    if err:
        return {"error": err}
    target, err = _validate_out_path(scope_root, out_path)
    if err:
        return {"error": err}

    # Write md to temp
    import tempfile
    with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False, encoding="utf-8") as tf:
        tf.write(md)
        md_path = tf.name

    target.parent.mkdir(parents=True, exist_ok=True)
    result = _run_cli([
        "pandoc", md_path,
        "-o", str(target),
        "--from", "markdown",
        "--to", target_format,
    ], timeout=120)
    try:
        Path(md_path).unlink()
    except Exception:
        pass

    if not result.get("ok"):
        return {"error": "pandoc failed", **result}
    return {
        "ok": True,
        "format": target_format,
        "path": out_path,
        "absolute": str(target),
        "size": target.stat().st_size if target.is_file() else 0,
    }


def tool_to_pdf(args: dict) -> dict:
    """libreoffice headless: any → PDF."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return {"error": "libreoffice not installed. brew install --cask libreoffice"}

    source_path = args.get("source_path", "")
    out_path = args.get("out_path")  # opzionale: stesso nome con .pdf
    scope_arg = args.get("scope")
    if not source_path:
        return {"error": "source_path required"}

    scope_root, err = _resolve_scope_root(scope_arg)
    if err:
        return {"error": err}
    source = (scope_root / source_path).resolve()
    try:
        source.relative_to(scope_root.resolve())
    except ValueError:
        return {"error": "source path outside scope"}
    if not source.is_file():
        return {"error": f"source file not found: {source_path}"}

    # Default out_path: stesso dir, stesso name ma .pdf
    if not out_path:
        out_path = str(Path(source_path).with_suffix(".pdf"))
    target, err = _validate_out_path(scope_root, out_path)
    if err:
        return {"error": err}

    target.parent.mkdir(parents=True, exist_ok=True)
    # soffice converte sempre nello stesso outdir del source. Usiamo tempdir e poi spostiamo.
    import tempfile
    with tempfile.TemporaryDirectory() as td:
        result = _run_cli([
            soffice, "--headless", "--convert-to", "pdf",
            "--outdir", td, str(source),
        ], timeout=180)
        if not result.get("ok"):
            return {"error": "libreoffice failed", **result}
        generated = Path(td) / (source.stem + ".pdf")
        if not generated.is_file():
            return {"error": "PDF not generated"}
        shutil.move(str(generated), str(target))

    return {
        "ok": True,
        "format": "pdf",
        "path": out_path,
        "absolute": str(target),
        "size": target.stat().st_size,
    }


def tool_from_markdown_slides(args: dict) -> dict:
    """marp-cli: markdown → slides pptx/pdf/html."""
    if not shutil.which("marp"):
        return {"error": "marp-cli not installed. npm install -g @marp-team/marp-cli"}

    md = args.get("md", "")
    theme = args.get("theme", "default")
    target_format = args.get("target_format", "pptx").lower()
    out_path = args.get("out_path", "")
    scope_arg = args.get("scope")
    if not md.strip():
        return {"error": "md required"}
    if target_format not in ("pptx", "pdf", "html"):
        return {"error": f"unsupported target_format: {target_format}"}

    scope_root, err = _resolve_scope_root(scope_arg)
    if err:
        return {"error": err}
    target, err = _validate_out_path(scope_root, out_path)
    if err:
        return {"error": err}

    import tempfile
    # Prepend marp frontmatter if missing
    if not md.lstrip().startswith("---"):
        md = f"---\nmarp: true\ntheme: {theme}\n---\n\n" + md

    with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False, encoding="utf-8") as tf:
        tf.write(md)
        md_path = tf.name

    target.parent.mkdir(parents=True, exist_ok=True)
    flag = {"pptx": "--pptx", "pdf": "--pdf", "html": "--html"}[target_format]
    result = _run_cli([
        "marp", flag, "--allow-local-files",
        "-o", str(target), md_path,
    ], timeout=180)
    try:
        Path(md_path).unlink()
    except Exception:
        pass
    if not result.get("ok"):
        return {"error": "marp failed", **result}
    return {
        "ok": True,
        "format": target_format,
        "theme": theme,
        "path": out_path,
        "absolute": str(target),
        "size": target.stat().st_size if target.is_file() else 0,
    }


def tool_deps_status(args: dict) -> dict:
    """Diagnostica deps installate."""
    status = {}
    # Python libs
    for lib in ("docx", "openpyxl", "pptx"):
        try:
            __import__(lib)
            status[lib] = True
        except ImportError:
            status[lib] = False
    # CLI binaries
    for bin_name in ("pandoc", "soffice", "libreoffice", "marp"):
        status[bin_name] = bool(shutil.which(bin_name))
    status["libreoffice_any"] = status.get("soffice") or status.get("libreoffice")
    status["py_office_complete"] = all([status["docx"], status["openpyxl"], status["pptx"]])
    return status


# ============================================================
# Tool registry
# ============================================================

TOOLS = [
    {
        "name": "office.generate_docx",
        "description": (
            "Genera un file Word (.docx) da spec JSON. Output salvato in scope_root/files/. "
            "Spec: {title, sections: [{heading, level, content, bullets?, table?, image?}]}. "
            "Table = list of lists. Image = path relativo (es. 'files/chart.png')."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "spec": {"type": "object", "description": "Document spec (title, sections)"},
                "out_path": {"type": "string", "description": "Path relativo (es. 'files/report-may.docx')"},
                "scope": {"type": "string", "description": "'hub' (default) o 'workspace:<name>'"},
            },
            "required": ["spec", "out_path"],
        },
    },
    {
        "name": "office.generate_xlsx",
        "description": (
            "Genera un file Excel (.xlsx) con sheet multipli, formule, charts. "
            "Spec: {sheets: [{name, rows: [[...], ...], charts?: [{type, data_range, title, anchor}]}]}. "
            "Type chart: bar|line|pie. data_range es. 'A1:B10'. anchor es. 'E2'."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "spec": {"type": "object", "description": "Workbook spec (sheets)"},
                "out_path": {"type": "string", "description": "Path relativo (es. 'files/data.xlsx')"},
                "scope": {"type": "string", "description": "'hub' o 'workspace:<name>'"},
            },
            "required": ["spec", "out_path"],
        },
    },
    {
        "name": "office.generate_pptx",
        "description": (
            "Genera presentazione PowerPoint (.pptx). "
            "Spec: {slides: [{layout: 'title'|'title_content'|'blank', title, bullets?, subtitle?, image?}]}. "
            "Per slides design avanzato → preferisci office.from_markdown_slides con marp."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "spec": {"type": "object", "description": "Presentation spec (slides)"},
                "out_path": {"type": "string", "description": "Path relativo (es. 'files/slides.pptx')"},
                "scope": {"type": "string"},
            },
            "required": ["spec", "out_path"],
        },
    },
    {
        "name": "office.from_markdown",
        "description": (
            "Converte markdown in docx/pptx/pdf/html via pandoc. "
            "Bidirezionale, supporta tables, footnotes, code blocks, metadata. "
            "Richiede `pandoc` installato (brew install pandoc)."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "md": {"type": "string", "description": "Markdown source"},
                "target_format": {"type": "string", "enum": ["docx", "pptx", "html", "pdf", "odt"]},
                "out_path": {"type": "string"},
                "scope": {"type": "string"},
            },
            "required": ["md", "target_format", "out_path"],
        },
    },
    {
        "name": "office.to_pdf",
        "description": (
            "Converte docx/xlsx/pptx/odt in PDF via LibreOffice headless. "
            "Richiede LibreOffice installato (brew install --cask libreoffice). "
            "out_path opzionale: default stesso nome con .pdf."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "source_path": {"type": "string", "description": "Path file source"},
                "out_path": {"type": "string", "description": "Opzionale, default <source>.pdf"},
                "scope": {"type": "string"},
            },
            "required": ["source_path"],
        },
    },
    {
        "name": "office.from_markdown_slides",
        "description": (
            "Markdown → slides (pptx/pdf/html) via marp-cli con temi (default/gaia/uncover). "
            "Richiede marp-cli (npm install -g @marp-team/marp-cli). "
            "Prepend frontmatter automatico se assente."
        ),
        "inputSchema": {
            "type": "object",
            "properties": {
                "md": {"type": "string"},
                "theme": {"type": "string", "description": "default|gaia|uncover"},
                "target_format": {"type": "string", "enum": ["pptx", "pdf", "html"]},
                "out_path": {"type": "string"},
                "scope": {"type": "string"},
            },
            "required": ["md", "out_path"],
        },
    },
    {
        "name": "office.deps_status",
        "description": "Diagnostica: ritorna quali deps Python e CLI sono installate (docx, openpyxl, pptx, pandoc, libreoffice, marp).",
        "inputSchema": {"type": "object", "properties": {}},
    },
]

TOOL_HANDLERS = {
    "office.generate_docx": tool_generate_docx,
    "office.generate_xlsx": tool_generate_xlsx,
    "office.generate_pptx": tool_generate_pptx,
    "office.from_markdown": tool_from_markdown,
    "office.to_pdf": tool_to_pdf,
    "office.from_markdown_slides": tool_from_markdown_slides,
    "office.deps_status": tool_deps_status,
}


def handle_request(req: dict):
    method = req.get("method")
    params = req.get("params") or {}
    req_id = req.get("id")
    if method == "initialize":
        return _ok(req_id, {
            "protocolVersion": PROTO_VERSION,
            "serverInfo": {"name": SERVER_NAME, "version": SERVER_VERSION},
            "capabilities": {"tools": {"listChanged": False}},
        })
    if method == "notifications/initialized":
        return None
    if method == "tools/list":
        return _ok(req_id, {"tools": TOOLS})
    if method == "tools/call":
        name = params.get("name")
        args = params.get("arguments") or {}
        handler = TOOL_HANDLERS.get(name)
        if not handler:
            return _err(req_id, -32601, f"unknown tool: {name}")
        try:
            result = handler(args)
            content = [{"type": "text", "text": json.dumps(result, ensure_ascii=False, indent=2)}]
            return _ok(req_id, {"content": content, "isError": "error" in result})
        except Exception as e:
            return _err(req_id, -32603, f"tool '{name}' failed: {type(e).__name__}: {e}")
    if method == "ping":
        return _ok(req_id, {})
    return _err(req_id, -32601, f"method not found: {method}")


def _ok(req_id, result):
    return {"jsonrpc": "2.0", "id": req_id, "result": result}


def _err(req_id, code, message, data=None):
    err = {"code": code, "message": message}
    if data is not None:
        err["data"] = data
    return {"jsonrpc": "2.0", "id": req_id, "error": err}


def main():
    print(f"[anja_office] starting (scope={SCOPE} root={ROOT})", file=sys.stderr, flush=True)
    for line in sys.stdin:
        line = line.strip()
        if not line:
            continue
        try:
            req = json.loads(line)
        except json.JSONDecodeError as e:
            err = _err(None, -32700, f"parse error: {e}")
            sys.stdout.write(json.dumps(err) + "\n")
            sys.stdout.flush()
            continue
        resp = handle_request(req)
        if resp is None:
            continue
        sys.stdout.write(json.dumps(resp, ensure_ascii=False) + "\n")
        sys.stdout.flush()


if __name__ == "__main__":
    main()
